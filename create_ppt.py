# import required things
from pptx import Presentation 
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
  
NEXT_BUTTON_URL = "https://cdn-icons-png.flaticon.com/512/3318/3318722.png"
PREVIOUS_BUTTON_URL = "https://image.freepik.com/free-icon/rewind-symbol_318-47556.jpg"

def set_text_box_title(slide, text, left, top, width, height):
    # Create white box with 0.7 opacity
    # Define dimensions and position
    # left = Cm(1.27)   # Distance from the left side of the slide
    # top = Cm(1.76)    # Distance from the top of the slide
    # width = Cm(22.86)   # Width of the text box
    # height = Cm(3.18)  # Height of the text box
    white_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = white_box.fill
    fill.solid()
    # fill.fore_color.rgb = RGBColor(255, 255, 0)  # white fill
    # fill.transparency = 0.3  # 30% transparency
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    # Add text on top of the white box
    # text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = text

    # Set text color
    font = text_frame.paragraphs[0].runs[0].font
    font.color.rgb = RGBColor(0, 0, 0)  # Black color

def set_slide_background(slide, image_url):
    slide_width = 9144000  # 10 inches
    slide_height = 6858000  # 7.5 inches

    # get image from URL
    image_stream = BytesIO(requests.get(image_url).content)
    image = Image.open(image_stream)
    
    # resize image
    image.thumbnail((slide_width, slide_height))

    # Convert PIL Image object to BytesIO object
    byte_img = BytesIO()
    image.save(byte_img, format='PNG')

    # Add image to slide
    slide.shapes.add_picture(byte_img, 0, 0, slide_width, slide_height)

def set_text_box_content(slide, text, left, top, width, height, r, g, b, size, isBold, br, bg, bb):
    # Create white box with 0.7 opacity
    # Define dimensions and position
    # left = Cm(1.27)   # Distance from the left side of the slide
    # top = Cm(5.01)    # Distance from the top of the slide
    # width = Cm(22.86)   # Width of the text box
    # height = Cm(12.57)  # Height of the text box
    white_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = white_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(br, bg, bb) 
    # fill.transparency = 0.3  # 30% transparency

    text_box = slide.shapes.add_textbox(left, top, width, height)
    # Add text on top of the white box
    # text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = text

    # Set text color
    font = text_frame.paragraphs[0].runs[0].font
    font.color.rgb = RGBColor(r, g, b)

    font.size = Pt(size)
    font.bold = isBold

def set_text_box_without_border(slide, text, left, top, width, height, r, g, b, size, isBold):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    for p in tf.paragraphs:
        font = p.runs[0].font
        font.size = Pt(size)
        font.color.rgb = RGBColor(r, g, b)
        font.name = 'Open Sans'
        font.bold = isBold

def add_filled_circle(slide, cx, cy, radius, fill_color):
    # add_filled_circle_cm(slide, 5, 5, 2, RGBColor(255, 0, 0))  # Adds a red filled circle.
    left = cx - radius
    top = cy - radius
    width = height = radius * 2

    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    
    fill = oval.fill
    fill.solid()
    fill.fore_color.rgb = fill_color

def draw_circle(slide, center_x, center_y, radius, line_color, line_width=1.0):

    left = center_x - radius
    top = center_y - radius
    width = height = 2 * radius
    
    # Add an oval (circle) shape to the slide
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    
    # Set the outline (line) properties of the circle
    line = circle.line
    line.color.rgb = line_color  # Set the line color
    line.width = line_width  # Set the line width
    
    # Remove the fill to make the circle transparent
    # fill = circle.fill
    # fill.solid()  # Set a solid fill
    # fill.fore_color.rgb = RGBColor(255, 255, 255)

def draw_line(slide, start_x, start_y, end_x, end_y, line_color, line_width=1.0):
    left = start_x
    top = start_y
    width = end_x - start_x
    height = end_y - start_y

    # Add a straight line shape to the slide
    line = slide.shapes.add_shape(MSO_SHAPE_TYPE.LINE, left, top, width, height)

    # Set the line properties
    line.line.color.rgb = line_color  # Set the line color
    line.line.width = Cm(line_width)  # Set the line width

def add_colored_chevron(slide, color, left, top, width, height):

    chevron = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    
    # Set the fill color of the chevron
    fill = chevron.fill
    fill.solid()
    fill.fore_color.rgb = color
    
    return chevron

def create_slide1(ppt):
    # To create blank slide layout
    # We have to use 6 as an argument
    # of slide_layouts  
    blank_slide_layout = ppt.slide_layouts[6] 
    
    # Attaching slide obj to slide
    slide = ppt.slides.add_slide(blank_slide_layout)
    
    # # For adjusting the  Margins in inches 
    # left = top = width = height = Inches(1) 
    
    # # creating textBox
    # txBox = slide.shapes.add_textbox(left, top,
    #                                 width, height)

    set_text_box_without_border(slide, "Presentation Skills", Cm(1.27), Cm(8.10), Cm(9.86), Cm(3.18), 0, 0, 128, 50, True)

    set_text_box_without_border(slide, "Duration: 35 minutes", Cm(1), Inches(8.5), Cm(7), Cm(3), 0, 0, 0, 18, True)

    left = Inches(7.8)
    top = Cm(0)
    slide.shapes.add_picture('slide1_pic1.jpg', left, top, height = Inches(9))

    set_text_box_content(slide, "Start", Inches(13.2), Inches(6.0), Cm(5), Cm(1.8), 255, 255, 255, 18, True, 0, 0, 255)
    set_text_box_content(slide, "Select Start to begin", Inches(13.2), Inches(7), Cm(6), Cm(1), 0, 0, 0, 18, True, 255, 255, 255)

def create_slide2(ppt):
    blank_slide_layout = ppt.slide_layouts[6] 
    slide = ppt.slides.add_slide(blank_slide_layout)
    set_text_box_without_border(slide, "Introduction", Inches(6), Cm(1.76), Cm(22.86), Cm(3.18), 0, 0, 128, 50, True)

    add_filled_circle(slide, Inches(8), Inches(4.5), Cm(4.4), RGBColor(232, 76, 56))
    add_filled_circle(slide, Inches(8), Inches(4.5), Cm(4.3), RGBColor(255, 255, 255))

    add_filled_circle(slide, Inches(5), Inches(4.5), Cm(2), RGBColor(255, 165, 0))
    add_filled_circle(slide, Inches(8), Inches(4.5), Cm(3), RGBColor(232, 76, 56))
    add_filled_circle(slide, Inches(11), Inches(4.5), Cm(2), RGBColor(173, 216, 230))

    slide.shapes.add_picture('slide2_pic1.png', Inches(4.6), Inches(4.3))
    slide.shapes.add_picture('slide2_pic1.png', Inches(7.6), Inches(4.3))
    slide.shapes.add_picture('slide2_pic1.png', Inches(10.6), Inches(4.3))

    slide.shapes.add_picture('slide2_pic2.png', Inches(3.7), Inches(3.2))
    slide.shapes.add_picture('slide2_pic3.png', Inches(9.7), Inches(3.2))

    draw_line(slide, Inches(5), Inches(4.5) + Cm(3.3), Inches(5), Inches(4.5) + Cm(5.3), RGBColor(255, 165, 0), 0.1)
    draw_line(slide, Inches(8), Inches(4.5) + Cm(4.2), Inches(8), Inches(4.5) + Cm(6.2), RGBColor(232, 76, 56), 0.1)
    draw_line(slide, Inches(11), Inches(4.5) + Cm(3.5), Inches(11), Inches(4.5) + Cm(5.5), RGBColor(173, 216, 230), 0.1)

    set_text_box_without_border(slide, "Infographic 1", Inches(4.5), Inches(4.5) + Cm(7), Cm(3), Cm(1), 0, 0, 0, 14, True)
    set_text_box_without_border(slide, "Infographic 2", Inches(7.5), Inches(4.5) + Cm(7), Cm(3), Cm(1), 0, 0, 0, 14, True)
    set_text_box_without_border(slide, "Infographic 3", Inches(10.5), Inches(4.5) + Cm(7), Cm(3), Cm(1), 0, 0, 0, 14, True)

    set_text_box_without_border(slide, "Morbi tincidunt ornare massa\neget egestas purus its viverra\nnullam accumsan", Inches(4.2), Inches(4.5) + Cm(8.1), Cm(3.3), Cm(2), 0, 0, 0, 11, False)
    set_text_box_without_border(slide, "Morbi tincidunt ornare massa\neget egestas purus its viverra\nnullam accumsan", Inches(7.2), Inches(4.5) + Cm(8.1), Cm(3.3), Cm(2), 0, 0, 0, 11, False)
    set_text_box_without_border(slide, "Morbi tincidunt ornare massa\neget egestas purus its viverra\nnullam accumsan", Inches(10.2), Inches(4.5) + Cm(8.1), Cm(3.3), Cm(2), 0, 0, 0, 11, False)

    response1 = requests.get(PREVIOUS_BUTTON_URL)
    if response1.status_code == 200:
        # Read the image content into a BytesIO object
        image_data = BytesIO(response1.content)
        slide.shapes.add_picture(image_data, Inches(15), Inches(8), height = Cm(0.7))

    response2 = requests.get(NEXT_BUTTON_URL)
    if response2.status_code == 200:
        # Read the image content into a BytesIO object
        image_data = BytesIO(response2.content)
        slide.shapes.add_picture(image_data, Inches(15.5), Inches(8), height = Cm(0.7))
    
def create_slide3(ppt):
    blank_slide_layout = ppt.slide_layouts[6] 
    slide = ppt.slides.add_slide(blank_slide_layout)
    set_text_box_without_border(slide, "Learning Objectives", Inches(5.2), Cm(1.76), Cm(22.86), Cm(3.18), 0, 0, 128, 50, True)

    slide.shapes.add_picture('slide3_pic1.png', Inches(6.5), Cm(5))

    add_colored_chevron(slide, RGBColor(232, 76, 56), Inches(2), Inches(5), Cm(4.9), Cm(1))
    add_colored_chevron(slide, RGBColor(255, 165, 0), Inches(2) + Cm(5), Inches(5), Cm(4.9), Cm(1))
    add_colored_chevron(slide, RGBColor(0, 255, 0), Inches(2) + Cm(10), Inches(5), Cm(4.9), Cm(1))
    add_colored_chevron(slide, RGBColor(173, 216, 230), Inches(2) + Cm(15), Inches(5), Cm(4.9), Cm(1))
    add_colored_chevron(slide, RGBColor(255, 105, 180), Inches(2) + Cm(20), Inches(5), Cm(4.9), Cm(1))

    add_filled_circle(slide, Inches(2) + Cm(2), Inches(5) - Cm(3), Cm(1.5), RGBColor(232, 76, 56))
    add_filled_circle(slide, Inches(2) + Cm(7), Inches(5) + Cm(4), Cm(1.5), RGBColor(255, 165, 0))
    add_filled_circle(slide, Inches(2) + Cm(12), Inches(5) - Cm(3), Cm(1.5), RGBColor(0, 255, 0))
    add_filled_circle(slide, Inches(2) + Cm(17), Inches(5) + Cm(4), Cm(1.5), RGBColor(173, 216, 230))
    add_filled_circle(slide, Inches(2) + Cm(22), Inches(5) - Cm(3), Cm(1.5), RGBColor(255, 105, 180))

    slide.shapes.add_picture('slide3_pic2.png', Inches(2) + Cm(2) - Cm(0.5), Inches(5) - Cm(3) - Cm(0.5))
    slide.shapes.add_picture('slide3_pic3.png', Inches(2) + Cm(7) - Cm(0.5), Inches(5) + Cm(4) - Cm(0.5))
    slide.shapes.add_picture('slide3_pic4.png', Inches(2) + Cm(12) - Cm(0.5), Inches(5) - Cm(3) - Cm(0.5))
    slide.shapes.add_picture('slide3_pic5.png', Inches(2) + Cm(17) - Cm(0.5), Inches(5) + Cm(4) - Cm(0.5))
    slide.shapes.add_picture('slide3_pic6.png', Inches(2) + Cm(22) - Cm(0.5), Inches(5) - Cm(3) - Cm(0.5))

    draw_line(slide, Inches(2) + Cm(2), Inches(5) - Cm(3) + Cm(1.4), Inches(2) + Cm(2), Inches(5) - Cm(3) + Cm(3.4), RGBColor(192, 192, 192), 0.1)
    draw_line(slide, Inches(2) + Cm(7), Inches(5) + Cm(4) - Cm(3.4), Inches(2) + Cm(7), Inches(5) + Cm(4) - Cm(1.4), RGBColor(192, 192, 192), 0.1)
    draw_line(slide, Inches(2) + Cm(12), Inches(5) - Cm(3) + Cm(1.4), Inches(2) + Cm(12), Inches(5) - Cm(3) + Cm(3.4), RGBColor(192, 192, 192), 0.1)
    draw_line(slide, Inches(2) + Cm(17), Inches(5) + Cm(4) - Cm(3.4), Inches(2) + Cm(17), Inches(5) + Cm(4) - Cm(1.4), RGBColor(192, 192, 192), 0.1)
    draw_line(slide, Inches(2) + Cm(22), Inches(5) - Cm(3) + Cm(1.4), Inches(2) + Cm(22), Inches(5) - Cm(3) + Cm(3.4), RGBColor(192, 192, 192), 0.1)

    set_text_box_without_border(slide, "Infographic 1", Inches(2) + Cm(2) - Cm(2), Inches(5) + Cm(2) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 14, True)
    set_text_box_without_border(slide, "Infographic 2", Inches(2) + Cm(7) - Cm(2), Inches(5) - Cm(3) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 14, True)
    set_text_box_without_border(slide, "Infographic 3", Inches(2) + Cm(12) - Cm(2), Inches(5) + Cm(2) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 14, True)
    set_text_box_without_border(slide, "Infographic 4", Inches(2) + Cm(17) - Cm(2), Inches(5) - Cm(3) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 14, True)
    set_text_box_without_border(slide, "Infographic 5", Inches(2) + Cm(22) - Cm(2), Inches(5) + Cm(2) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 14, True)

    set_text_box_without_border(slide, "Morbi tincidunt ornar eacu\nmassa eget egestas\npurus its viverra nullam", Inches(2) + Cm(2) - Cm(2), Inches(5) + Cm(3) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 11, False)
    set_text_box_without_border(slide, "Morbi tincidunt ornar eacu\nmassa eget egestas\npurus its viverra nullam", Inches(2) + Cm(7) - Cm(2), Inches(5) - Cm(2) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 11, False)
    set_text_box_without_border(slide, "Morbi tincidunt ornar eacu\nmassa eget egestas\npurus its viverra nullam", Inches(2) + Cm(12) - Cm(2), Inches(5) + Cm(3) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 11, False)
    set_text_box_without_border(slide, "Morbi tincidunt ornar eacu\nmassa eget egestas\npurus its viverra nullam", Inches(2) + Cm(17) - Cm(2), Inches(5) - Cm(2) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 11, False)
    set_text_box_without_border(slide, "Morbi tincidunt ornar eacu\nmassa eget egestas\npurus its viverra nullam", Inches(2) + Cm(22) - Cm(2), Inches(5) + Cm(3) - Cm(0.5), Cm(22.86), Cm(3.18), 0, 0, 128, 11, False)

    response1 = requests.get(PREVIOUS_BUTTON_URL)
    if response1.status_code == 200:
        # Read the image content into a BytesIO object
        image_data = BytesIO(response1.content)
        slide.shapes.add_picture(image_data, Inches(15), Inches(8), height = Cm(0.7))

    response2 = requests.get(NEXT_BUTTON_URL)
    if response2.status_code == 200:
        # Read the image content into a BytesIO object
        image_data = BytesIO(response2.content)
        slide.shapes.add_picture(image_data, Inches(15.5), Inches(8), height = Cm(0.7))

def create_slide4(ppt):
    blank_slide_layout = ppt.slide_layouts[6] 
    slide = ppt.slides.add_slide(blank_slide_layout)
    set_text_box_without_border(slide, "The Four Types of\nPresentation", Inches(1), Cm(4.76), Cm(22.86), Cm(3.18), 0, 0, 128, 50, True)

    slide.shapes.add_picture('slide3_pic1.png', Inches(1.2), Cm(9.76))

    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8)-Cm(0.3), Cm(4)-Cm(0.3), Cm(15.9), Inches(7)+Cm(0.6))
    
    slide.shapes.add_picture('slide4_pic1.jpg', Inches(8), Cm(4), height = Inches(7))

    set_text_box_without_border(slide, "There are four kinds of presentations:\ninformational, instructional,\nstimulating, and convincing.", Inches(1.2), Cm(11.76), Cm(22.86), Cm(3.18), 0, 0, 0, 18, False)

    response1 = requests.get(PREVIOUS_BUTTON_URL)
    if response1.status_code == 200:
        # Read the image content into a BytesIO object
        image_data = BytesIO(response1.content)
        slide.shapes.add_picture(image_data, Inches(15), Inches(8), height = Cm(0.7))

    response2 = requests.get(NEXT_BUTTON_URL)
    if response2.status_code == 200:
        # Read the image content into a BytesIO object
        image_data = BytesIO(response2.content)
        slide.shapes.add_picture(image_data, Inches(15.5), Inches(8), height = Cm(0.7))

def create_slide5(ppt):
    blank_slide_layout = ppt.slide_layouts[6] 
    slide = ppt.slides.add_slide(blank_slide_layout)
    set_text_box_without_border(slide, "About Us", Inches(1), Cm(3.76), Cm(22.86), Cm(3.18), 232, 76, 56, 15, False)
    set_text_box_without_border(slide, "Technical Skills", Inches(1), Cm(4.76), Cm(22.86), Cm(3.18), 0, 0, 128, 50, True)

    slide.shapes.add_picture('slide3_pic1.png', Inches(1.2), Cm(7.76))

    box1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Cm(4), Cm(2), Cm(2))
    fill = box1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(232, 76, 56)

    box1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5) + Cm(8), Cm(4), Cm(2), Cm(2))
    fill = box1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(232, 76, 56)

    box1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Cm(4)+Cm(8), Cm(2), Cm(2))
    fill = box1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(232, 76, 56)

    box1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5)+Cm(8), Cm(4)+Cm(8), Cm(2), Cm(2))
    fill = box1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(232, 76, 56)

    slide.shapes.add_picture('slide5_pic1.png', Inches(8.7), Cm(4.6))
    slide.shapes.add_picture('slide5_pic2.png', Inches(8.7)+ Cm(8), Cm(4.6))
    slide.shapes.add_picture('slide5_pic3.png', Inches(8.7), Cm(4.6)+ Cm(8))
    slide.shapes.add_picture('slide5_pic4.png', Inches(8.7)+ Cm(8), Cm(4.6)+ Cm(8))

    response1 = requests.get(PREVIOUS_BUTTON_URL)
    if response1.status_code == 200:
        # Read the image content into a BytesIO object
        image_data = BytesIO(response1.content)
        slide.shapes.add_picture(image_data, Inches(15), Inches(8), height = Cm(0.7))

    response2 = requests.get(NEXT_BUTTON_URL)
    if response2.status_code == 200:
        # Read the image content into a BytesIO object
        image_data = BytesIO(response2.content)
        slide.shapes.add_picture(image_data, Inches(15.5), Inches(8), height = Cm(0.7))



def create_ppt():
    # Creating Object
    ppt = Presentation() 
    ppt.slide_width = Inches(16)
    ppt.slide_height = Inches(9)

    create_slide1(ppt)
    create_slide2(ppt)
    create_slide3(ppt)
    create_slide4(ppt)
    create_slide5(ppt)
    
    # # creating textFrames
    # tf = txBox.text_frame
    # tf.text = "This is text inside a textbox"
    
    # # adding Paragraphs
    # p = tf.add_paragraph() 
    
    # # adding text
    # p.text = "This is a second paragraph that's bold and italic" 
    
    # # font 
    # p.font.bold = True
    # p.font.italic = True
    
    # p = tf.add_paragraph()
    # p.text = "This is a third paragraph that's big " 
    # p.font.size = Pt(40)
    
    # save file
    ppt.save('test_2.pptx')
    
    print("done")

if __name__ == "__main__":
    create_ppt()